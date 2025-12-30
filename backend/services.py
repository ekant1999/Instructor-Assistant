from __future__ import annotations

import asyncio
import json
import os
import re
import uuid
from pathlib import Path
from tempfile import NamedTemporaryFile
from typing import Any, AsyncGenerator, Dict, List, Optional, Tuple
import logging

import requests
from litellm import acompletion, completion
from pptx import Presentation
from pypdf import PdfReader

from backend.core.database import get_conn
from backend.core.questions import render_canvas_markdown
from .mcp_client import MCPClientError, call_tool as call_mcp_tool, is_configured as mcp_configured

logging.basicConfig(level=os.getenv("LOG_LEVEL", "INFO").upper())
logger = logging.getLogger(__name__)
from .schemas import (
    PaperChatMessage,
    Question,
    QuestionContextUploadResponse,
    QuestionGenerationRequest,
    QuestionGenerationResponse,
    QuestionInsertionRequest,
)

# DEFAULT_MODEL = os.getenv("LITELLM_MODEL", "gpt-4o-mini")
DEFAULT_MODEL = os.getenv("LITELLM_MODEL", "gpt-5-mini")
_raw_temp = os.getenv("LITELLM_TEMPERATURE")
if _raw_temp is not None:
    TEMPERATURE = float(_raw_temp)
elif "gpt-5" in DEFAULT_MODEL.lower():
    TEMPERATURE = 1.0
else:
    TEMPERATURE = 0.2
MAX_TOKENS = int(os.getenv("LITELLM_MAX_TOKENS", "4000"))
MAX_CONTEXT_CHARS = int(os.getenv("QUESTION_CONTEXT_CHAR_LIMIT", "60000"))
DEFAULT_PROVIDER = (os.getenv("LLM_PROVIDER") or "openai").strip().lower()
SUMMARY_PROVIDER = (os.getenv("SUMMARY_PROVIDER") or os.getenv("LLM_PROVIDER") or "openai").strip().lower()
LOCAL_LLM_URL = (os.getenv("LOCAL_LLM_URL") or "http://localhost:11434").rstrip("/")
LOCAL_LLM_MODEL = os.getenv("LOCAL_LLM_MODEL", "llama3.1:8b")
LOCAL_LLM_TIMEOUT = int(os.getenv("LOCAL_LLM_TIMEOUT", "60"))
LOCAL_TOOL_LIMIT = int(os.getenv("LOCAL_TOOL_STEPS", "12"))
ALLOWED_LOCAL_TOOLS = {"list_contexts", "read_context"}
TYPE_PATTERNS = {
    "mcq": r"mcqs?|multiple\s+choice(?:\s+questions?)?",
    "short_answer": r"short[-\s]?answer(?:\s+questions?)?",
    "true_false": r"true\s*(?:or|/)?\s*false(?:\s+questions?)?|tf",
    "essay": r"essay(?:\s+questions?)?",
}


class QuestionGenerationError(RuntimeError):
    """Raised when the LLM output cannot be parsed."""


def _completion_limit_args(model_name: str) -> Dict[str, Any]:
    """Return the correct token-limit argument for the selected model."""
    if "gpt-5" in (model_name or "").lower():
        return {"max_completion_tokens": MAX_TOKENS}
    return {"max_tokens": MAX_TOKENS}


def generate_questions(payload: QuestionGenerationRequest) -> QuestionGenerationResponse:
    provider = _resolve_provider(payload)
    messages = _build_messages(payload, provider)
    if provider == "local":
        return _generate_questions_local(payload, messages)
    return _generate_questions_openai(payload, messages)


def _generate_questions_openai(payload: QuestionGenerationRequest, messages: List[Dict[str, str]]) -> QuestionGenerationResponse:
    if not os.getenv("OPENAI_API_KEY") and not os.getenv("LITELLM_API_KEY"):
        raise QuestionGenerationError("OPENAI_API_KEY (or LITELLM_API_KEY) must be set to use the question generator.")

    comp_kwargs = _completion_limit_args(DEFAULT_MODEL)
    try:
        response = completion(
            model=DEFAULT_MODEL,
            messages=messages,
            temperature=TEMPERATURE,
            **comp_kwargs,
        )
    except Exception as exc:
        raise QuestionGenerationError(f"LLM request failed: {exc}") from exc

    content = response["choices"][0]["message"]["content"]
    questions = _parse_questions(content)
    markdown = render_canvas_markdown(payload.instructions, [q.model_dump() for q in questions], {})

    return QuestionGenerationResponse(
        questions=questions,
        markdown=markdown,
        raw_response=content,
    )


async def stream_generate_questions(payload: QuestionGenerationRequest) -> AsyncGenerator[Dict[str, Any], None]:
    provider = _resolve_provider(payload)
    if provider == "local":
        result = _generate_questions_local(payload, _build_messages(payload, provider))
        yield {
            "type": "complete",
            "questions": [q.model_dump() for q in result.questions],
            "markdown": result.markdown,
            "raw_response": result.raw_response,
        }
        return

    if not os.getenv("OPENAI_API_KEY") and not os.getenv("LITELLM_API_KEY"):
        raise QuestionGenerationError("OPENAI_API_KEY (or LITELLM_API_KEY) must be set to use the question generator.")

    messages = _build_messages(payload)
    comp_kwargs = _completion_limit_args(DEFAULT_MODEL)
    try:
        stream = await acompletion(
            model=DEFAULT_MODEL,
            messages=messages,
            temperature=TEMPERATURE,
            stream=True,
            **comp_kwargs,
        )
    except Exception as exc:
        raise QuestionGenerationError(f"LLM request failed: {exc}") from exc

    collected: List[str] = []
    async for chunk in stream:
        delta = chunk["choices"][0]["delta"]
        content = delta.get("content")
        if not content:
            continue
        collected.append(content)
        yield {"type": "chunk", "content": content}

    full_text = "".join(collected)
    questions = _parse_questions(full_text)
    markdown = render_canvas_markdown(payload.instructions, [q.model_dump() for q in questions], {})
    yield {
        "type": "complete",
        "questions": [q.model_dump() for q in questions],
        "markdown": markdown,
        "raw_response": full_text,
    }


def _list_available_contexts() -> List[Dict[str, Any]]:
    if not mcp_configured():
        raise QuestionGenerationError("LOCAL_MCP_SERVER_URL must be configured to list contexts.")
    try:
        payload = call_mcp_tool("list_contexts", {})
    except MCPClientError as exc:
        raise QuestionGenerationError(f"Failed to list contexts via MCP: {exc}") from exc
    contexts = payload.get("contexts") if isinstance(payload, dict) else None
    if isinstance(contexts, list):
        return contexts
    return []


def _build_messages(payload: QuestionGenerationRequest, provider: str = "openai") -> List[Dict[str, str]]:
    derived_type_counts, derived_total = _derive_type_counts(payload.instructions)
    type_instruction = "Feel free to use MCQ, short_answer, true_false, or essay question types."
    if payload.question_types:
        type_instruction = (
            f"The instructor prefers the following question types (use these labels when possible): {', '.join(payload.question_types)}."
        )
    elif derived_type_counts:
        summary = ", ".join(f"{count} {kind.replace('_', ' ')}" for kind, count in derived_type_counts.items())
        type_instruction = (
            f"You must generate exactly these question counts/types: {summary}. Use only the labels mcq, true_false, short_answer, or essay."
        )

    total_questions = payload.question_count or derived_total
    count_instruction = "Generate only the requested questions."
    if total_questions is not None:
        count_instruction = f"Generate exactly {total_questions} questions."
    if derived_type_counts and derived_total:
        count_instruction = f"Generate exactly {derived_total} questions total, matching the per-type counts above."

    context_block = f"\nContext:\n{payload.context.strip()}" if payload.context else ""
    mcq_rule = (
        "For every multiple choice question, provide four distinct answer options (a, b, c, d) in the order presented, "
        "unless the instructions specify a different number of options."
    )

    schema_block = """
Return JSON with this shape:
{
  "questions": [
    {
      "kind": "mcq | short_answer | true_false | essay",
      "text": "Question text",
      "options": ["optional list of options"],
      "answer": "short answer or letter",
      "explanation": "why the answer is correct",
      "reference": "source citation"
    }
  ]
}
"""

    system_prompt = (
        "You are an experienced instructor who writes exam-ready questions. "
        "Only produce valid JSON. Avoid commentary outside JSON."
    )

    constraint_note = (
        "If you cannot satisfy the requested counts/types, respond with a JSON error object like {\"error\": \"reason\"} instead of returning questions."
    )
    if derived_type_counts:
        constraint_note += " Do not output question types that were not explicitly requested."

    user_prompt = (
        f"{payload.instructions.strip()}\n\n"
        f"{type_instruction}\n"
        f"{count_instruction}\n"
        f"{schema_block}\n"
        f"{mcq_rule}\n"
        f"{context_block}\n"
        f"Ensure answers reflect the context. {constraint_note}"
    )

    if provider == "local":
        return _build_local_messages(payload, user_prompt)

    return [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_prompt},
    ]


def _build_local_messages(payload: QuestionGenerationRequest, base_prompt: str) -> List[Dict[str, str]]:
    contexts = _list_available_contexts()
    if contexts:
        context_lines = "\n".join(
            f"- {ctx.get('context_id')}: {ctx.get('filename')} ({ctx.get('characters')} chars, preview: {(ctx.get('preview') or '')[:120]}...)"
            for ctx in contexts
        )
        context_note = (
            "Uploaded files are available through tools. Call 'list_contexts' to see all IDs, then 'read_context' to fetch excerpts.\n"
            f"Known files:\n{context_lines}\n"
        )
    else:
        context_note = (
            "No uploaded documents are currently available. If you still need context, ask the instructor for more details."
        )

    tool_instructions = """
You must respond ONLY with JSON objects. Valid patterns:
- {"action":"call_tool","tool":"list_contexts","arguments":{}}
- {"action":"call_tool","tool":"read_context","arguments":{"context_id":"CTX_ID","start":0,"length":4000}}
- {"action":"final","content":<questions JSON matching the schema>}

Never send plain text or Markdown. If your previous reply wasn't valid JSON, immediately send a JSON reminder or the correct response.

Tools:
- list_contexts: returns metadata for each uploaded file (context_id, filename, preview, length).
- read_context: arguments {context_id (string, required), start (int, optional), length (int <= 4000, optional)}. Returns a text excerpt.

When you have enough information, produce the final JSON with the required "questions" array and use {"action":"final", "content": ...}.

Important rules (violating any of these requires you to immediately stop and output {"action":"final","content":{"error":"CONSTRAINT_VIOLATION"}}):
- You must NEVER call any tool other than list_contexts or read_context. Do not invent new tool names under any circumstance.
- Every read_context call must include a valid context_id you previously saw in list_contexts. Never send read_context without context_id or with an unknown ID.
- If a tool call fails or you realize you cannot comply with these instructions, stop immediately and return the CONSTRAINT_VIOLATION error above.
"""

    system_prompt = (
        "You are an experienced instructor who can read uploaded source materials via tools and then write exam-ready questions. "
        "Do not invent file contents—call the tools when you need information. "
        "Follow the JSON-only protocol strictly."
    )

    user_prompt = f"{base_prompt}\n\n{context_note}\n{tool_instructions}"
    return [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_prompt},
    ]


def _parse_questions(content: str) -> List[Question]:
    cleaned = _strip_code_fences(content)
    payload: Any
    try:
        payload = json.loads(cleaned)
    except json.JSONDecodeError as exc:
        raise QuestionGenerationError(f"LLM response was not valid JSON: {exc}") from exc

    raw_questions = payload.get("questions") if isinstance(payload, dict) else payload
    if not isinstance(raw_questions, list):
        raise QuestionGenerationError("LLM response did not include a 'questions' list.")

    normalized: List[Question] = []
    for entry in raw_questions:
        if not isinstance(entry, dict):
            continue
        text = entry.get("text") or entry.get("question")
        if not text:
            continue
        options = entry.get("options")
        if isinstance(options, list):
            options = [str(o).strip() for o in options if str(o).strip()]
        else:
            options = None
        answer_raw = entry.get("answer") or entry.get("solution") or ""
        explanation_raw = entry.get("explanation") or entry.get("rationale") or ""
        reference_raw = entry.get("reference") or entry.get("source") or ""
        question = Question(
            kind=(entry.get("kind") or entry.get("type") or "short_answer").lower(),
            text=text.strip(),
            options=options,
            answer=str(answer_raw).strip() or None,
            explanation=str(explanation_raw).strip() or None,
            reference=str(reference_raw).strip() or None,
        )
        normalized.append(question)

    if not normalized:
        raise QuestionGenerationError("LLM response did not include any valid questions.")

    return normalized


def _strip_code_fences(raw: str) -> str:
    pattern = re.compile(r"^```(?:json)?\s*(.*?)\s*```$", re.DOTALL)
    match = pattern.match(raw.strip())
    if match:
        return match.group(1)
    return raw


def _derive_type_counts(instructions: str) -> Tuple[Dict[str, int], Optional[int]]:
    counts: Dict[str, int] = {}
    text = instructions.lower()

    for kind, pattern in TYPE_PATTERNS.items():
        regex = re.compile(rf"(\d+)\s*(?:{pattern})", re.I)
        for match in regex.finditer(text):
            qty = int(match.group(1))
            counts[kind] = counts.get(kind, 0) + qty

    total = sum(counts.values()) if counts else None
    if total is None:
        general = re.search(r"(\d+)\s+(?:questions|items)", text)
        if general:
            total = int(general.group(1))

    return counts, total


def _generate_questions_local(payload: QuestionGenerationRequest, messages: List[Dict[str, str]]) -> QuestionGenerationResponse:
    raw = _run_local_tool_session(messages)
    questions = _parse_questions(raw)
    markdown = render_canvas_markdown(payload.instructions, [q.model_dump() for q in questions], {})
    return QuestionGenerationResponse(
        questions=questions,
        markdown=markdown,
        raw_response=raw,
    )


def _call_local_llm(messages: List[Dict[str, str]]) -> str:
    if not LOCAL_LLM_URL:
        raise QuestionGenerationError("Set LOCAL_LLM_URL to use the local LLM provider.")
    url = f"{LOCAL_LLM_URL}/api/chat"
    payload = {
        "model": LOCAL_LLM_MODEL,
        "messages": messages,
        "stream": False,
    }
    try:
        response = requests.post(url, json=payload, timeout=LOCAL_LLM_TIMEOUT)
    except requests.RequestException as exc:
        raise QuestionGenerationError(f"Local LLM request failed: {exc}") from exc
    if response.status_code >= 400:
        raise QuestionGenerationError(f"Local LLM error: {response.text}")
    data = response.json()
    content = None
    if isinstance(data.get("message"), dict):
        content = data["message"].get("content")
    if not content and isinstance(data.get("messages"), list) and data["messages"]:
        content = data["messages"][-1].get("content")
    if not content and isinstance(data.get("response"), str):
        content = data["response"]
    if not content:
        raise QuestionGenerationError("Local LLM did not return any content.")
    return content


def _run_local_tool_session(messages: List[Dict[str, str]]) -> str:
    conversation = [dict(msg) for msg in messages]
    for step in range(LOCAL_TOOL_LIMIT):
        reply = _call_local_llm(conversation)
        cleaned = reply.strip()
        data = _try_parse_json(cleaned)
        if isinstance(data, dict) and data.get("action") == "call_tool":
            tool_name = data.get("tool")
            if tool_name not in ALLOWED_LOCAL_TOOLS:
                raise QuestionGenerationError(
                    f"Local LLM attempted unsupported tool '{tool_name}'. Only list_contexts/read_context are allowed."
                )
            arguments = data.get("arguments") or {}
            logger.info("[local-llm] tool=%s args=%s", tool_name, arguments)
            tool_output = _execute_local_tool(tool_name, arguments)
            conversation.append({"role": "assistant", "content": cleaned})
            conversation.append(
                {
                    "role": "user",
                    "content": (
                        f"Tool '{tool_name}' result:\n{json.dumps(tool_output)}\n"
                        "Continue following the instructions. Remember to reply with JSON when calling tools or when producing the final answer."
                    ),
                }
            )
            continue
        if isinstance(data, dict) and data.get("action") == "final":
            final_content = data.get("content")
            if not final_content:
                raise QuestionGenerationError("Local LLM returned an empty final payload.")
            if isinstance(final_content, (dict, list)):
                final_content = json.dumps(final_content)
            else:
                final_content = str(final_content)
            logger.info("[local-llm] final response after %s steps", step + 1)
            return final_content
        logger.warning("[local-llm] unstructured response: %s", cleaned[:200])
        conversation.append({"role": "assistant", "content": cleaned})
        conversation.append(
            {
                "role": "user",
                "content": (
                    "Reminder: respond only with JSON objects. Use {\"action\":\"call_tool\",...} for tool calls or {\"action\":\"final\",...} when returning the final questions."
                ),
            }
        )
        continue
    raise QuestionGenerationError("Local LLM exceeded tool call limit before returning a valid final response.")


def _execute_local_tool(name: Optional[str], arguments: Dict[str, Any]) -> Dict[str, Any]:
    if not mcp_configured():
        raise QuestionGenerationError("LOCAL_MCP_SERVER_URL must be configured to call tools.")
    try:
        return call_mcp_tool(name or "", arguments or {})
    except MCPClientError as exc:
        raise QuestionGenerationError(f"MCP tool '{name}' failed: {exc}") from exc


def _try_parse_json(raw: str) -> Optional[Any]:
    try:
        return json.loads(_strip_code_fences(raw))
    except (json.JSONDecodeError, TypeError):
        return None


def _resolve_provider(payload: QuestionGenerationRequest) -> str:
    provider = (payload.provider or DEFAULT_PROVIDER or "openai").strip().lower()
    if provider not in {"openai", "local"}:
        provider = DEFAULT_PROVIDER or "openai"
    return provider


def generate_insertion_preview(
    question_set_payload: Dict[str, Any],
    payload: QuestionInsertionRequest,
) -> Tuple[List[Question], List[Question], int]:
    existing = question_set_payload.get("questions") or []
    question_models = [_question_from_dict(item) for item in existing]
    anchor_index, anchor_label = _resolve_insert_index(question_models, payload.anchor_question_id, payload.position)
    anchor_sentence = _build_anchor_instruction(anchor_label, payload.position)
    user_instruction = payload.instructions.strip()
    combined_instruction = f"{user_instruction}\n\n{anchor_sentence}\nReturn only the new questions that should be inserted at that location."
    request = QuestionGenerationRequest(
        instructions=combined_instruction,
        context=payload.context,
        question_count=payload.question_count,
        question_types=payload.question_types,
        provider=payload.provider,
    )
    result = generate_questions(request)
    preview_questions = result.questions
    merged_questions = (
        question_models[:anchor_index] + preview_questions + question_models[anchor_index:]
    )
    return preview_questions, merged_questions, anchor_index


def _question_from_dict(data: Dict[str, Any]) -> Question:
    try:
        return Question.model_validate(data)
    except Exception:
        return Question(
            id=data.get("id"),
            kind=(data.get("kind") or "short_answer"),
            text=data.get("text") or "Untitled question",
            options=data.get("options"),
            answer=data.get("answer"),
            explanation=data.get("explanation"),
            reference=data.get("reference"),
        )


def _resolve_insert_index(questions: List[Question], anchor_id: Optional[int], position: str) -> Tuple[int, str]:
    if not questions:
        return 0, "the beginning of the question set"
    if anchor_id is None:
        if position == "before":
            return 0, f"before question '{questions[0].text[:80]}'"
        return len(questions), f"after question '{questions[-1].text[:80]}'"
    for idx, question in enumerate(questions):
        if question.id == anchor_id:
            if position == "before":
                return idx, f"before question '{question.text[:80]}'"
            return idx + 1, f"after question '{question.text[:80]}'"
    # Fallback to end if anchor not found
    return len(questions), f"after question '{questions[-1].text[:80]}'"


def _build_anchor_instruction(anchor_phrase: str, position: str) -> str:
    return (
        f"Insert the new questions {anchor_phrase}. "
        "Do not repeat existing questions; output only the additional ones."
    )


def summarize_paper_chat(paper_id: int, messages: List[PaperChatMessage]) -> Dict[str, Any]:
    with get_conn() as conn:
        paper = conn.execute("SELECT id, title FROM papers WHERE id=?", (paper_id,)).fetchone()
        if not paper:
            raise QuestionGenerationError("Paper not found.")
        sections = conn.execute(
            "SELECT page_no, text FROM sections WHERE paper_id=? ORDER BY page_no ASC",
            (paper_id,)
        ).fetchall()
    context = "\n\n".join((row["text"] or "" for row in sections))[:MAX_CONTEXT_CHARS]
    if not context.strip():
        raise QuestionGenerationError("No text available for this paper.")

    system_prompt = (
        "You are a research assistant. Summarize the given paper and answer follow-up questions using only the provided context."
    )
    base_messages: List[Dict[str, str]] = [
        {"role": "system", "content": f"{system_prompt}\nContext:\n{context}"}
    ]
    for msg in messages:
        base_messages.append({"role": msg.role, "content": msg.content})

    if SUMMARY_PROVIDER == "local":
        try:
            text = _call_local_llm(base_messages).strip()
        except Exception as exc:
            raise QuestionGenerationError(f"Local LLM request failed: {exc}") from exc
    else:
        try:
            response = completion(
                model=DEFAULT_MODEL,
                messages=base_messages,
                temperature=TEMPERATURE,
                **_completion_limit_args(DEFAULT_MODEL),
            )
        except Exception as exc:
            raise QuestionGenerationError(f"LLM request failed: {exc}") from exc

        text = response["choices"][0]["message"]["content"].strip()
    note_title = paper["title"] or "Paper Summary"
    return {
        "message": text,
        "paper_id": paper["id"],
        "paper_title": paper["title"],
        "suggested_title": note_title,
    }


async def extract_context_from_upload(filename: str, data: bytes) -> QuestionContextUploadResponse:
    suffix = Path(filename or "").suffix.lower()
    if suffix not in {".pdf", ".ppt", ".pptx"}:
        raise QuestionGenerationError("Only PDF and PPT/PPTX files are supported at the moment.")

    loop = asyncio.get_running_loop()
    text = await loop.run_in_executor(None, _extract_text_from_bytes, filename, data)
    text = text.strip()
    if not text:
        raise QuestionGenerationError("Could not extract any text from the uploaded file.")
    if len(text) > MAX_CONTEXT_CHARS:
        text = text[:MAX_CONTEXT_CHARS]
    preview = text[:400].strip()
    return QuestionContextUploadResponse(
        context_id=uuid.uuid4().hex,
        filename=filename or "upload",
        characters=len(text),
        preview=preview,
        text=text,
    )


def _extract_text_from_bytes(filename: str, data: bytes) -> str:
    suffix = Path(filename or "").suffix.lower()
    with NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(data)
        tmp_path = Path(tmp.name)
    try:
        if suffix == ".pdf":
            return _extract_pdf_text(tmp_path)
        return _extract_ppt_text(tmp_path)
    finally:
        tmp_path.unlink(missing_ok=True)


def _extract_pdf_text(path: Path) -> str:
    reader = PdfReader(str(path))
    parts: List[str] = []
    for page in reader.pages:
        parts.append(page.extract_text() or "")
    return "\n".join(parts)


def _extract_ppt_text(path: Path) -> str:
    presentation = Presentation(str(path))
    parts: List[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                parts.append(shape.text)
    return "\n".join(parts)
